import os
import platform
from fpdf import FPDF
from PIL import Image
from docx2pdf import convert as docx_convert
from pptx import Presentation
from reportlab.pdfgen import canvas
from reportlab.lib.pagesizes import A4

ASCII_ART = r"""
██████╗ ██████╗ ███╗   ██╗██╗   ██╗███████╗██████╗ ████████╗███████╗██████╗ 
██╔════╝██╔═══██╗████╗  ██║██║   ██║██╔════╝██╔══██╗╚══██╔══╝██╔════╝██╔══██╗
██║     ██║   ██║██╔██╗ ██║██║   ██║█████╗  ██████╔╝   ██║   █████╗  ██████╔╝
██║     ██║   ██║██║╚██╗██║╚██╗ ██╔╝██╔══╝  ██╔══██╗   ██║   ██╔══╝  ██╔══██╗
╚██████╗╚██████╔╝██║ ╚████║ ╚████╔╝ ███████╗██║  ██║   ██║   ███████╗██║  ██║
 ╚═════╝ ╚═════╝ ╚═╝  ╚═══╝  ╚═══╝  ╚══════╝╚═╝  ╚═╝   ╚═╝   ╚══════╝╚═╝  ╚═╝
"""

def clear_terminal():
    os.system('cls' if platform.system() == 'Windows' else 'clear')
    print(ASCII_ART)

desktop = os.path.join(os.path.expanduser('~'), 'Desktop')

def convert_txt_to_pdf(file_path, save_path):
    print("\nQuantas linhas por célula você deseja no PDF?")
    lines_per_cell = input("Escolha (1/2/3): ")
    if lines_per_cell not in ['1','2','3']:
        lines_per_cell = '1'
    lines_per_cell = int(lines_per_cell)

    pdf = FPDF()
    pdf.add_page()
    pdf.set_auto_page_break(auto=True, margin=15)
    pdf.set_font("Arial", size=12)

    with open(file_path, 'r', encoding='utf-8') as f:
        buffer = []
        for line in f:
            buffer.append(line.rstrip('\n'))
            if len(buffer) == lines_per_cell:
                pdf.multi_cell(0, 10, ' | '.join(buffer))
                buffer = []
        if buffer:
            pdf.multi_cell(0, 10, ' | '.join(buffer))
    pdf.output(save_path)

def convert_image_to_pdf(file_path, save_path):
    image = Image.open(file_path).convert('RGB')
    image.save(save_path)

def convert_docx_to_pdf(file_path, save_path):
    docx_convert(file_path, save_path)

def convert_pptx_to_pdf(file_path, save_path):
    prs = Presentation(file_path)
    c = canvas.Canvas(save_path, pagesize=A4)
    width, height = A4
    for slide in prs.slides:
        c.setFont("Helvetica", 14)
        y_offset = height - 50
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip() != "":
                text = shape.text.strip()
                c.drawString(50, y_offset, text)
                y_offset -= 20
        c.showPage()
    c.save()

def convert_to_pdf(file_path):
    filename = os.path.splitext(os.path.basename(file_path))[0] + ".pdf"
    save_path = os.path.join(desktop, filename)
    ext = os.path.splitext(file_path)[1].lower()

    try:
        if ext in ['.jpg', '.jpeg', '.png', '.bmp', '.tiff']:
            convert_image_to_pdf(file_path, save_path)
        elif ext == '.txt':
            convert_txt_to_pdf(file_path, save_path)
        elif ext == '.docx':
            convert_docx_to_pdf(file_path, save_path)
        elif ext == '.pptx':
            convert_pptx_to_pdf(file_path, save_path)
        elif ext == '.pdf':
            print("Já é um PDF!")
            return
        else:
            print(f"Formato {ext} não suportado.")
            return
        print(f"\nArquivo convertido com sucesso: {save_path}")
    except Exception as e:
        print("Erro ao converter:", e)

def search_files():
    search_name = input("Digite parte do nome do arquivo: ").lower()
    all_files = []
    common_dirs = [
        os.path.join(os.path.expanduser('~'), 'Desktop'),
        os.path.join(os.path.expanduser('~'), 'Documents'),
        os.path.join(os.path.expanduser('~'), 'Downloads')
    ]
    for folder in common_dirs:
        if os.path.exists(folder):
            for root, dirs, files in os.walk(folder):
                for file in files:
                    if search_name in file.lower():
                        all_files.append(os.path.join(root, file))

    if not all_files:
        print("Nenhum arquivo encontrado!")
        input("\nPressione Enter para continuar...")
        clear_terminal()
        return

    print("\nArquivos encontrados:")
    for i, f in enumerate(all_files, 1):
        print(f"{i} | {os.path.basename(f)}")

    choice = input("\nQual deseja converter? Digite o número: ")
    if not choice.isdigit() or int(choice) < 1 or int(choice) > len(all_files):
        print("Número inválido!")
        input("\nPressione Enter para continuar...")
        clear_terminal()
        return

    convert_to_pdf(all_files[int(choice)-1])
    input("\nPressione Enter para continuar...")
    clear_terminal()

def direct_convert():
    file_path = input("Digite o caminho completo do arquivo: ")
    if not os.path.exists(file_path):
        print("Arquivo não encontrado!")
        input("\nPressione Enter para continuar...")
        clear_terminal()
        return
    convert_to_pdf(file_path)
    input("\nPressione Enter para continuar...")
    clear_terminal()

def main():
    clear_terminal()
    while True:
        print("\nMenu:")
        print("1 | Search e converter")
        print("2 | Converter pelo caminho do arquivo")
        print("3 | Sair")
        choice = input("Escolha uma opção: ")
        if choice == '1':
            clear_terminal()
            search_files()
        elif choice == '2':
            clear_terminal()
            direct_convert()
        elif choice == '3':
            break
        else:
            print("Opção inválida!")
            input("\nPressione Enter para continuar...")
            clear_terminal()

if __name__ == "__main__":
    main()